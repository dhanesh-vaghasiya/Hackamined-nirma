from pptx import Presentation
from pptx.util import Inches

prs = Presentation('document/devxlabs_skills_mirage_deck.pptx')
for i, slide in enumerate(prs.slides):
    print(f"=== SLIDE {i+1} ===")
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            print(shape.text.strip())
        if shape.has_table:
            table = shape.table
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                print(" | ".join(cells))
    print()
